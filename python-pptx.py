import warnings
warnings.filterwarnings('ignore')
from pptx import Presentation
from nltk.corpus import wordnet
from nltk.stem import WordNetLemmatizer
import nltk 
from nltk.corpus import stopwords 
from nltk.tokenize import word_tokenize, sent_tokenize 
from autocorrect import spell
from mic import voice_assitant
import win32com.client
from lemmatizer import preprocessing

app = win32com.client.Dispatch("PowerPoint.Application")

presentation = app.Presentations.Open(FileName=u"C:\\Users\\ejimhos\\Downloads\\CV Presentation.pptx")

def return_tags_dict(prs):
	tags = {}
	count = 0
	for slide_no, slide in enumerate(prs.slides):
		notes_slide = slide.notes_slide
		data = notes_slide.notes_text_frame.text
		if len(data)==0:
			continue
		tags[slide_no] = []
		for word in data.split(','):
			tags[slide_no].append(preprocessing(word)[0])
	return tags

def check(word,tags,present_slide):
	try:
		if word in set(tags[present_slide+1]):
			return (True,'next')
		elif word in set(tags[present_slide-1]):
			return (True,'previous')
		else:
			return (False,-1)
	except:
		return (False,-1)

def check_list(text,tags,present_slide):
	for word in text:
		out = check(word,tags,present_slide)
		print(out)
		if out[0]:
			return out[1]
	return -1

if __name__ == "__main__":
	print("[INFO] Loading Presentation ......\n")
	prs = Presentation("CV Presentation.pptx")
	tags = return_tags_dict(prs)
	print(tags,str('\n'))
	presentation.SlideShowSettings.Run()
	count = 0
	print("[INFO] Assitant Unlocked ......\n")
	while(1):
		text = list(set(preprocessing(voice_assitant())))
		# text = preprocessing(input("Enter the Command: "))
		print(f'[INFO] You Said {text}\n')
		if len(text) == 1:
			if text[0].lower() == 'next':
				presentation.SlideShowWindow.View.Next()
				count+=1
				print(f'[INFO] Present Slide is {count}')
			elif text[0].lower() == 'previous':
				presentation.SlideShowWindow.View.Previous()
				if count!=0:
					count-=1
				else:
					print("[INFO]Presentation on first slide......\n")
				print(f'[INFO] Present Slide is {count}\n')
			elif text[0].lower() == 'close':
				presentation.SlideShowWindow.View.Exit()
				break
			else:
				slide_info = check_list(text,tags,count)
				if slide_info != -1:
					if slide_info == 'next':
						presentation.SlideShowWindow.View.Next()
						count+=1
						print(f'[INFO] Present Slide is {count}')
					elif slide_info == 'previous':
						presentation.SlideShowWindow.View.Previous()
						if count!=0:
							count-=1
						else:
							print("[INFO]Presentation on first slide......\n")
						print(f'[INFO] Present Slide is {count}\n')
				else:
					continue
		else:
			slide_info = check_list(text, tags, count)
			if slide_info != -1:
				if slide_info == 'next':
					presentation.SlideShowWindow.View.Next()
					count+=1
					print(f'[INFO] Present Slide is {count}\n')
				elif slide_info == 'previous':
					presentation.SlideShowWindow.View.Previous()
					if count!=0:
						count-=1
					else:
						print('[INFO] Presentation on first slide ......\n')
					print(f'[INFO] Present Slide is {count}\n')
			else:
				continue
	app.Quit()


